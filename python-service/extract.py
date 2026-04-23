from pptx import Presentation
import sys

file_path = sys.argv[1]

prs = Presentation(file_path)

text = ""

for slide in prs.slides:
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text += shape.text + "\n"

print(text)
